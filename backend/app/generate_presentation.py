import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


TITLE_FONT = "Segoe UI"
BODY_FONT = "Segoe UI"


def add_title_slide(prs: Presentation, title: str, subtitle: str, footer: str | None = None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    if footer:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(9), Inches(0.4))
        p = txBox.text_frame.paragraphs[0]
        p.text = footer
        p.font.size = Pt(14)
        p.font.name = BODY_FONT
        p.alignment = PP_ALIGN.CENTER


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str]):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = BODY_FONT
        p.font.size = Pt(24)


def add_section_slide(prs: Presentation, title: str, sections: list[tuple[str, list[str]]]):
    # Title and Content, but we will build rich content box
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    left = Inches(0.8)
    top = Inches(1.8)
    width = Inches(8.8)
    height = Inches(4.8)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()

    for idx, (section_title, items) in enumerate(sections):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = section_title
        p.level = 0
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.name = TITLE_FONT
        for item in items:
            q = tf.add_paragraph()
            q.text = item
            q.level = 1
            q.font.size = Pt(22)
            q.font.name = BODY_FONT


def build_presentation(output_path: str):
    prs = Presentation()

    # Theme-like defaults via placeholders
    for layout in prs.slide_layouts:
        pass

    # Slide 1
    add_title_slide(
        prs,
        title="Customer Support Chatbot",
        subtitle="Realtime Chat, Tickets, and FAQs in One Place",
        footer="Faster, clearer customer support",
    )

    # Slide 2
    add_bullets_slide(
        prs,
        title="The Problem — Traditional Support Pain",
        bullets=[
            "Long waits and slow responses",
            "Same questions asked again and again",
            "High effort for simple, repeat issues",
            "Inconsistent experiences across tools",
            "No simple path to 24/7 coverage",
        ],
    )

    # Slide 3
    add_section_slide(
        prs,
        title="Our Solution — What The App Provides",
        sections=[
            ("Realtime Chat (Socket.IO)", [
                "Customers chat instantly",
                "First message auto-creates a ticket",
            ]),
            ("Ticket Management (Admin Pages)", [
                "Track status: Open, In Progress, Resolved, Closed",
                "Full chat history attached to each ticket",
            ]),
            ("FAQ Knowledge Base", [
                "Manage common Q&A",
                "Help customers self-serve faster",
            ]),
            ("Optional AI Replies", [
                "Enable with OpenAI/Google keys",
            ]),
        ],
    )

    # Slide 4
    add_section_slide(
        prs,
        title="Customer Experience — The Chat Flow",
        sections=[
            ("Simple Start", [
                "Enter name, email, subject, category",
            ]),
            ("Live Conversation", [
                "Realtime chat with instant updates",
            ]),
            ("Ticket Created Automatically", [
                "Every chat is tracked from the first message",
            ]),
            ("Faster Answers", [
                "Customers can also check FAQs",
            ]),
        ],
    )

    # Slide 5
    add_section_slide(
        prs,
        title="Admin Dashboard — Control and Clarity",
        sections=[
            ("Live Case View", [
                "See all tickets with statuses",
            ]),
            ("Filters & Search", [
                "Filter by status/category; search by subject/customer",
            ]),
            ("Full Context", [
                "Ticket + complete chat history",
            ]),
            ("Manage Knowledge", [
                "Create and edit FAQs",
            ]),
        ],
    )

    # Slide 6
    add_section_slide(
        prs,
        title="FAQ Management — How We Improve Over Time",
        sections=[
            ("Collect", ["Identify frequent questions from resolved tickets"]),
            ("Curate", ["Add or update FAQs in Admin"]),
            ("Categorize", ["Group by topic for easy discovery"]),
            ("Reuse", ["Use FAQs to resolve faster"]),
        ],
    )

    # Slide 7
    add_section_slide(
        prs,
        title="Technical Architecture — What It’s Built On",
        sections=[
            ("Frontend", [
                "React + Vite + Tailwind (fast UI)",
            ]),
            ("Backend", [
                "FastAPI REST + Socket.IO (ASGI) for realtime chat",
            ]),
            ("Data Stores", [
                "PostgreSQL: users, tickets, FAQs",
                "MongoDB: chat transcripts/case memory",
                "Redis: activity timestamps and flags",
            ]),
            ("Platform", [
                "Docker Compose local setup; JWT auth; CORS enabled",
            ]),
            ("Key Behaviors", [
                "Auto-create ticket on first chat message",
                "Ticket lifecycle: Open → In Progress → Resolved → Closed",
                "Optional AI connectors (OpenAI/Google)",
            ]),
        ],
    )

    output_dir = os.path.dirname(output_path)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)
    prs.save(output_path)


if __name__ == "__main__":
    build_presentation(os.path.join(os.getcwd(), "CustomerSupportChatbot.pptx"))


